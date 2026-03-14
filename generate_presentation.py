#!/usr/bin/env python3
"""
MoveMint_Business_Case.pptx generator
======================================
Senior Product Design + Business Strategy · 10-slide deck
Handshake AI application · Palette: Deep Space Crypto
Fonts: Calibri Light (headlines) / Calibri (body)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import subprocess, os

# ══════════════════════════════════════════════════════════════════════════════
#  DESIGN SYSTEM
# ══════════════════════════════════════════════════════════════════════════════

C = {
    "bg":      RGBColor(0x0D, 0x11, 0x17),   # near-black canvas
    "card":    RGBColor(0x16, 0x1B, 0x22),   # elevated card surface
    "card2":   RGBColor(0x1C, 0x24, 0x2E),   # secondary card
    "blue":    RGBColor(0x58, 0xA6, 0xFF),   # electric blue
    "purple":  RGBColor(0xBC, 0x8C, 0xFF),   # soft purple
    "green":   RGBColor(0x3F, 0xB9, 0x50),   # emerald
    "amber":   RGBColor(0xF7, 0x85, 0x16),   # amber / warning
    "red":     RGBColor(0xFF, 0x6B, 0x6B),   # red (negative)
    "white":   RGBColor(0xFF, 0xFF, 0xFF),
    "light":   RGBColor(0xE6, 0xED, 0xF3),   # primary text
    "muted":   RGBColor(0x8B, 0x94, 0x9E),   # secondary text
    "divider": RGBColor(0x21, 0x29, 0x32),   # subtle rule
    "ghost":   RGBColor(0x1A, 0x21, 0x2B),   # decorative bg numbers
}

FH = "Calibri Light"   # headline font
FB = "Calibri"         # body font
FM = "Courier New"     # mono font

W = Inches(13.33)
H = Inches(7.5)

# ══════════════════════════════════════════════════════════════════════════════
#  PRIMITIVE HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def prs_new():
    p = Presentation()
    p.slide_width = W
    p.slide_height = H
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

# ── shapes ──────────────────────────────────────────────────────────────────

def rect(slide, x, y, w, h, fill=None, lc=None, lw=None):
    """Rectangle. lc=line color, lw=line width Pt."""
    sh = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if lc:
        sh.line.color.rgb = lc
        sh.line.width = Pt(lw or 1)
    else:
        # make line same colour as background → invisible
        sh.line.color.rgb = fill or C["bg"]
        sh.line.width = Pt(0.1)
    return sh

def bg(slide, color):
    rect(slide, 0, 0, W, H, fill=color)

# ── text box factory ─────────────────────────────────────────────────────────

def txt(slide, text, x, y, w, h,
        fn=FB, fs=16, bold=False, italic=False,
        color=None, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.name = fn
    r.font.size = Pt(fs)
    r.font.bold = bold
    r.font.italic = italic
    if color:
        r.font.color.rgb = color
    return tb

def mtxt(slide, lines, x, y, w, h,
         align=PP_ALIGN.LEFT, wrap=True):
    """
    Multi-paragraph textbox.
    lines = [{"t": "...", "fn": FH, "fs": 28, "bold": True, "color": C["white"],
               "italic": False, "sa": 6, "align": PP_ALIGN.LEFT}]
    sa = space_after in Pt
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, d in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = d.get("align", align)
        if d.get("sa"):
            p.space_after = Pt(d["sa"])
        r = p.add_run()
        r.text = d.get("t", "")
        r.font.name = d.get("fn", FB)
        r.font.size = Pt(d.get("fs", 16))
        r.font.bold = d.get("bold", False)
        r.font.italic = d.get("italic", False)
        if d.get("color"):
            r.font.color.rgb = d["color"]
    return tb

# ── reusable composite elements ──────────────────────────────────────────────

def slide_chrome(slide, num: str, section: str):
    """Shared chrome: left accent bar, top-right accent, slide label."""
    rect(slide, 0,          0, Inches(0.055), H,           fill=C["blue"])
    rect(slide, W - Inches(3), 0, Inches(3),  Inches(0.06), fill=C["purple"])
    txt(slide, f"{num} — {section.upper()}",
        Inches(0.18), Inches(0.18), Inches(8), Inches(0.28),
        fn=FB, fs=7.5, bold=True, color=C["muted"])

def section_label(slide, label, x, y, color=None):
    txt(slide, label.upper(),
        x, y, Inches(4), Inches(0.28),
        fn=FB, fs=8, bold=True, color=color or C["blue"])

def kpi_card(slide, big, label, x, y, w=Inches(2.8), h=Inches(1.5),
             fill=C["card"], big_color=None, lc=None):
    """Big-number KPI card."""
    rect(slide, x, y, w, h, fill=fill, lc=lc or fill, lw=0.1)
    txt(slide, big,
        x + Inches(0.18), y + Inches(0.12), w - Inches(0.3), Inches(0.85),
        fn=FH, fs=38, bold=True, color=big_color or C["blue"], wrap=False)
    txt(slide, label,
        x + Inches(0.18), y + Inches(0.92), w - Inches(0.3), Inches(0.5),
        fn=FB, fs=11, color=C["muted"])

def pill(slide, text, x, y, fill=None, text_color=None):
    """Small rounded-rect pill label (faked with rect)."""
    w = Inches(len(text) * 0.085 + 0.3)
    rect(slide, x, y, w, Inches(0.27), fill=fill or C["blue"])
    txt(slide, text, x + Inches(0.1), y + Inches(0.03),
        w - Inches(0.15), Inches(0.24),
        fn=FB, fs=9, bold=True, color=text_color or C["bg"])

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

# ─── SLIDE 1 · THE HOOK ──────────────────────────────────────────────────────
def s01_hook(prs):
    s = blank(prs)
    bg(s, C["bg"])
    slide_chrome(s, "01", "The Opportunity")

    # Ghost grid lines (subtle texture)
    for i in range(1, 7):
        rect(s, 0, Inches(i * 1.1), W, Inches(0.012), fill=C["divider"])

    # Decorative big ghost number top-right
    txt(s, "∞", Inches(9.5), Inches(0.5), Inches(3.5), Inches(5.5),
        fn=FH, fs=320, bold=True, color=C["ghost"], wrap=False)

    # Main headline block
    txt(s, "Everyone wants",
        Inches(0.55), Inches(1.0), Inches(9.5), Inches(1.25),
        fn=FH, fs=68, bold=True, color=C["white"])
    txt(s, "to launch a token.",
        Inches(0.55), Inches(2.05), Inches(9.5), Inches(1.25),
        fn=FH, fs=68, bold=True, color=C["white"])

    # Accent divider
    rect(s, Inches(0.55), Inches(3.4), Inches(4.8), Inches(0.05), fill=C["blue"])

    txt(s, "Almost no one can.",
        Inches(0.55), Inches(3.55), Inches(9.0), Inches(1.1),
        fn=FH, fs=44, bold=True, color=C["blue"])

    txt(s, "The technical barrier to deploying, trading, and graduating a token\n"
           "to a decentralized exchange has locked billions in latent creator value.",
        Inches(0.55), Inches(4.75), Inches(7.5), Inches(0.95),
        fn=FB, fs=16, color=C["muted"])

    # Bottom wordmark
    txt(s, "MOVEMINT", Inches(9.4), Inches(6.5), Inches(3.6), Inches(0.7),
        fn=FH, fs=34, bold=True, color=C["purple"], align=PP_ALIGN.RIGHT)
    txt(s, "Token launcher · Aptos blockchain",
        Inches(9.4), Inches(7.05), Inches(3.6), Inches(0.35),
        fn=FB, fs=10, color=C["muted"], align=PP_ALIGN.RIGHT)
    return s


# ─── SLIDE 2 · THE GAP ───────────────────────────────────────────────────────
def s02_gap(prs):
    s = blank(prs)
    bg(s, C["bg"])
    slide_chrome(s, "02", "The Market Gap")

    # Split: left panel (55 %) · right panel (45 %)
    rect(s, 0, 0, Inches(7.5), H, fill=C["card"])
    rect(s, Inches(7.5), 0, Inches(5.83), H, fill=C["bg"])

    # Left — Problems
    section_label(s, "Why current solutions fail", Inches(0.55), Inches(0.55))
    txt(s, "The Broken Landscape",
        Inches(0.55), Inches(0.9), Inches(6.5), Inches(0.75),
        fn=FH, fs=30, bold=True, color=C["white"])

    problems = [
        ("Pump.fun dominance",    "Solana-only — Aptos creators have zero equivalent product."),
        ("High technical barrier","Deploying a Move smart contract requires weeks of expertise."),
        ("No liquidity path",     "DIY tokens die on launch — no automated graduation to DEX."),
        ("Opaque price discovery","No on-chain bonding curve means prices are manipulable."),
        ("Fragmented UX",         "Creators juggle 4+ tools: wallet, indexer, DEX, explorer."),
    ]
    y = Inches(1.85)
    for title, body in problems:
        rect(s, Inches(0.55), y, Inches(0.06), Inches(0.55), fill=C["red"])
        txt(s, title, Inches(0.75), y, Inches(6.3), Inches(0.35),
            fn=FB, fs=13, bold=True, color=C["light"])
        txt(s, body,  Inches(0.75), y + Inches(0.33), Inches(6.3), Inches(0.32),
            fn=FB, fs=11, color=C["muted"])
        y += Inches(0.9)

    # Right — The shift
    section_label(s, "The Macro Tailwind", Inches(7.85), Inches(0.55), C["green"])
    txt(s, "Why Now",
        Inches(7.85), Inches(0.9), Inches(5.0), Inches(0.75),
        fn=FH, fs=30, bold=True, color=C["white"])

    stats = [
        ("$847M",  "Aptos TVL — fastest-growing L1 in 2024"),
        ("4.2M+",  "monthly active wallets on Aptos"),
        ("0",      "dedicated token launchers on Aptos today"),
        ("$130B",  "memecoin market cap at peak — demand is proven"),
    ]
    y = Inches(1.85)
    for big, label in stats:
        kpi_card(s, big, label, Inches(7.85), y,
                 w=Inches(4.9), h=Inches(1.05),
                 fill=C["card2"], big_color=C["green"])
        y += Inches(1.22)

    # Vertical divider line
    rect(s, Inches(7.47), Inches(0.5), Inches(0.04), Inches(6.5), fill=C["blue"])
    return s


# ─── SLIDE 3 · THE THESIS ────────────────────────────────────────────────────
def s03_thesis(prs):
    s = blank(prs)
    bg(s, C["bg"])
    slide_chrome(s, "03", "Investment Thesis")

    # Top accent band
    rect(s, 0, Inches(0.55), W, Inches(0.55), fill=C["card"])

    section_label(s, "The core thesis", Inches(0.55), Inches(0.62))

    # Central thesis statement — large
    txt(s, "Aptos is the next home for the token creator economy.",
        Inches(0.55), Inches(1.3), Inches(12.2), Inches(1.1),
        fn=FH, fs=38, bold=True, color=C["white"])

    txt(s, "MoveMint is the one product that makes that possible — today.",
        Inches(0.55), Inches(2.25), Inches(12.2), Inches(0.7),
        fn=FH, fs=24, bold=False, color=C["blue"])

    # 3-column evidence pillars
    pillars = [
        (C["blue"],   "SUPPLY",
         "Zero competitors on Aptos",
         "No pump.fun equivalent exists on Aptos. MoveMint captures a completely open market with no incumbent to displace."),
        (C["purple"], "DEMAND",
         "Creators want permissionless tools",
         "Memecoins + micro-cap tokens generated $130B in peak market cap. Demand is structurally proven — just not on Aptos."),
        (C["green"],  "INFRASTRUCTURE",
         "Aptos is production-ready",
         "Sub-second finality, $847M TVL, Move language safety. The L1 is ready. The app layer is missing."),
    ]
    col_x = Inches(0.4)
    col_w = Inches(4.0)
    for accent, tag, title, body in pillars:
        rect(s, col_x, Inches(3.3), col_w, Inches(3.6), fill=C["card"])
        rect(s, col_x, Inches(3.3), col_w, Inches(0.08), fill=accent)
        section_label(s, tag, col_x + Inches(0.25), Inches(3.5), accent)
        txt(s, title,
            col_x + Inches(0.25), Inches(3.85), col_w - Inches(0.45), Inches(0.65),
            fn=FH, fs=20, bold=True, color=C["white"])
        txt(s, body,
            col_x + Inches(0.25), Inches(4.55), col_w - Inches(0.45), Inches(2.2),
            fn=FB, fs=13, color=C["muted"])
        col_x += Inches(4.45)

    # Footer rule
    rect(s, Inches(0.55), Inches(7.18), Inches(12.2), Inches(0.04), fill=C["divider"])
    txt(s, "A defensible first-mover on Aptos has near-zero time to claim.",
        Inches(0.55), Inches(7.22), Inches(12.0), Inches(0.25),
        fn=FB, fs=10, color=C["muted"])
    return s


# ─── SLIDE 4 · THE INNOVATION ────────────────────────────────────────────────
def s04_product(prs):
    s = blank(prs)
    bg(s, C["bg"])
    slide_chrome(s, "04", "The Innovation")

    # Hero left panel
    rect(s, 0, 0, Inches(5.2), H, fill=C["card"])

    txt(s, "MOVEMINT",
        Inches(0.45), Inches(1.0), Inches(4.5), Inches(0.9),
        fn=FH, fs=52, bold=True, color=C["white"])

    pill(s, "LIVE ON APTOS TESTNET", Inches(0.45), Inches(1.95),
         fill=C["green"], text_color=C["bg"])

    txt(s, "A permissionless platform to create,\ntrade, and graduate tokens on Aptos —\nin under 60 seconds.",
        Inches(0.45), Inches(2.35), Inches(4.35), Inches(1.5),
        fn=FB, fs=16, color=C["light"])

    # KPIs
    mini_kpis = [("0.2 APT", "launch fee"), ("1%", "trade fee"), ("1,283 APT", "graduation target")]
    ky = Inches(4.1)
    for big, lbl in mini_kpis:
        txt(s, big, Inches(0.45), ky, Inches(2.2), Inches(0.5),
            fn=FH, fs=26, bold=True, color=C["blue"])
        txt(s, lbl, Inches(0.45), ky + Inches(0.45), Inches(2.2), Inches(0.3),
            fn=FB, fs=10, color=C["muted"])
        ky += Inches(0.85)

    # Right — 6 feature cards (2×3 grid)
    features = [
        (C["blue"],   "Token Factory",    "Deploy any fungible token in one transaction via Move smart contract."),
        (C["purple"], "Bonding Curve",    "Automated price discovery: buy pressure drives price up, sell returns APT."),
        (C["green"],  "DEX Graduation",   "Tokens hitting 1,283 APT auto-migrate to Hyperion DEX with locked LP."),
        (C["amber"],  "Real-time Feed",   "Live trade history, price ticks, and leaderboard via Geomi GraphQL index."),
        (C["blue"],   "Multi-wallet",     "Petra, Fewcha, Martian, Rise, MiZu — plus Google OAuth."),
        (C["purple"], "Watchlist",        "Persistent, per-user watchlist stored in localStorage for quick access."),
    ]
    gx = Inches(5.45)
    gy = Inches(0.6)
    cw = Inches(3.8)
    ch = Inches(2.15)
    for i, (accent, title, body) in enumerate(features):
        col = i % 2
        row = i // 2
        cx = gx + col * Inches(4.0)
        cy = gy + row * Inches(2.28)
        rect(s, cx, cy, cw, ch, fill=C["card2"])
        rect(s, cx, cy, Inches(0.07), ch, fill=accent)
        txt(s, title,
            cx + Inches(0.22), cy + Inches(0.2), cw - Inches(0.32), Inches(0.48),
            fn=FH, fs=15, bold=True, color=C["light"])
        txt(s, body,
            cx + Inches(0.22), cy + Inches(0.7), cw - Inches(0.32), Inches(1.3),
            fn=FB, fs=12, color=C["muted"])

    rect(s, Inches(5.2), 0, Inches(0.04), H, fill=C["divider"])
    return s


# ─── SLIDE 5 · TECHNICAL DEEP DIVE ──────────────────────────────────────────
def s05_tech(prs):
    s = blank(prs)
    bg(s, C["bg"])
    slide_chrome(s, "05", "Technical Architecture")

    # Subtle background grid
    for i in range(1, 8):
        rect(s, Inches(i * 1.7), 0, Inches(0.012), H, fill=C["divider"])

    txt(s, "How It Works Under the Hood",
        Inches(0.55), Inches(0.55), Inches(12.0), Inches(0.7),
        fn=FH, fs=32, bold=True, color=C["white"])

    # 3-layer architecture
    layers = [
        {
            "accent": C["blue"],
            "tag":    "LAYER 1 · SMART CONTRACTS",
            "title":  "Move Contracts",
            "items": [
                "1,089-line Move module on Aptos",
                "Bonding curve formula:",
                "  price = 19,029,514,756 /",
                "  (800M − tokens_sold) + 61.9",
                "Events: Created, Purchase,",
                "Sale, Graduated, Debug",
                "Pre-grad fees: 0.9% platform",
                "+ 0.1% creator",
            ],
        },
        {
            "accent": C["purple"],
            "tag":    "LAYER 2 · DATA LAYER",
            "title":  "Indexer + RPC",
            "items": [
                "Geomi No-Code Indexer (GraphQL)",
                "Aptos Fullnode RPC for balances",
                "10-min TTL cache + dedup",
                "Page-visibility polling pause",
                "Custom hooks:",
                "  useTokenData, useGraduation,",
                "  useTokenMarketData,",
                "  useAptPrice",
            ],
        },
        {
            "accent": C["green"],
            "tag":    "LAYER 3 · DEX INTEGRATION",
            "title":  "Hyperion Graduation",
            "items": [
                "Hyperion SDK full-range pool",
                "Threshold: 1,283 APT raised",
                "→ 1,200 APT migrated to DEX",
                "LP tokens burned (dead addr)",
                "Post-grad fees: 0.05% platform",
                "0.2% creator · 0.05% LP",
                "GraduationListener component",
                "monitors events continuously",
            ],
        },
    ]

    bx = Inches(0.4)
    bw = Inches(4.1)
    for lyr in layers:
        rect(s, bx, Inches(1.4), bw, Inches(5.7), fill=C["card"])
        rect(s, bx, Inches(1.4), bw, Inches(0.06), fill=lyr["accent"])
        section_label(s, lyr["tag"], bx + Inches(0.22), Inches(1.58), lyr["accent"])
        txt(s, lyr["title"],
            bx + Inches(0.22), Inches(1.92), bw - Inches(0.35), Inches(0.6),
            fn=FH, fs=22, bold=True, color=C["white"])
        iy = Inches(2.62)
        for item in lyr["items"]:
            ic = C["muted"] if item.startswith("  ") else C["light"]
            txt(s, item.strip(),
                bx + Inches(0.22 if not item.startswith("  ") else 0.42),
                iy, bw - Inches(0.5), Inches(0.38),
                fn=FM if item.startswith("  ") else FB,
                fs=11, color=ic)
            iy += Inches(0.42)
        bx += Inches(4.45)

    # Arrow connectors between layers
    for ax in [Inches(4.52), Inches(8.97)]:
        txt(s, "→", ax, Inches(4.0), Inches(0.4), Inches(0.5),
            fn=FH, fs=24, bold=True, color=C["blue"], align=PP_ALIGN.CENTER)

    # React/TS badge
    txt(s, "Frontend: React 18 · TypeScript · @aptos-labs/ts-sdk v1.39.0",
        Inches(0.55), Inches(7.15), Inches(12.0), Inches(0.3),
        fn=FM, fs=10, color=C["muted"])
    return s


# ─── SLIDE 6 · USER JOURNEY ──────────────────────────────────────────────────
def s06_journey(prs):
    s = blank(prs)
    bg(s, C["bg"])
    slide_chrome(s, "06", "User Journey")

    txt(s, "From Idea to DEX Listing in 5 Steps",
        Inches(0.55), Inches(0.55), Inches(12.0), Inches(0.7),
        fn=FH, fs=32, bold=True, color=C["white"])
    txt(s, "Total elapsed time for a creator: < 60 seconds to launch · < 1 hour to graduate (market-dependent)",
        Inches(0.55), Inches(1.25), Inches(12.0), Inches(0.38),
        fn=FB, fs=13, color=C["muted"])

    # Horizontal timeline steps
    steps = [
        (C["blue"],   "01", "Connect Wallet",
         "Petra / Google OAuth", "< 10s"),
        (C["purple"], "02", "Create Token",
         "Name, symbol, supply,\nimage metadata", "0.2 APT fee"),
        (C["amber"],  "03", "Community Buys",
         "Bonding curve drives\nprice discovery", "Auto"),
        (C["green"],  "04", "Graduation",
         "1,283 APT raised →\nHyperion pool created", "83 APT fee"),
        (C["blue"],   "05", "DEX Trading",
         "Full AMM liquidity,\nLP permanently locked", "∞"),
    ]

    sw  = Inches(2.35)
    sh  = Inches(4.5)
    sx  = Inches(0.35)
    sy  = Inches(2.0)

    for i, (accent, num, title, detail, time_) in enumerate(steps):
        # Card
        rect(s, sx, sy, sw, sh, fill=C["card"])
        rect(s, sx, sy, sw, Inches(0.08), fill=accent)

        # Step number (big ghost)
        txt(s, num, sx + Inches(0.1), sy + Inches(0.1), sw - Inches(0.1), Inches(1.1),
            fn=FH, fs=72, bold=True, color=C["ghost"], wrap=False)

        txt(s, title,
            sx + Inches(0.2), sy + Inches(1.05), sw - Inches(0.3), Inches(0.55),
            fn=FH, fs=17, bold=True, color=C["white"])
        txt(s, detail,
            sx + Inches(0.2), sy + Inches(1.65), sw - Inches(0.3), Inches(1.5),
            fn=FB, fs=12, color=C["muted"])

        # Time badge
        rect(s, sx + Inches(0.2), sy + sh - Inches(0.65),
             Inches(1.6), Inches(0.38), fill=accent)
        txt(s, time_,
            sx + Inches(0.3), sy + sh - Inches(0.62),
            Inches(1.4), Inches(0.35),
            fn=FB, fs=11, bold=True, color=C["bg"])

        # Arrow (not after last)
        if i < len(steps) - 1:
            txt(s, "›", sx + sw + Inches(0.08), sy + Inches(1.8),
                Inches(0.3), Inches(0.5),
                fn=FH, fs=28, bold=True, color=C["blue"])

        sx += sw + Inches(0.3)

    # Bottom note
    rect(s, Inches(0.35), Inches(6.8), Inches(12.6), Inches(0.55), fill=C["card2"])
    txt(s, "⚡  GraduationListener runs in the background — creators never have to manually trigger the DEX migration.",
        Inches(0.55), Inches(6.87), Inches(12.2), Inches(0.4),
        fn=FB, fs=12, color=C["light"])
    return s


# ─── SLIDE 7 · MARKET STRATEGY & REVENUE ─────────────────────────────────────
def s07_revenue(prs):
    s = blank(prs)
    bg(s, C["bg"])
    slide_chrome(s, "07", "Market Strategy & Revenue")

    txt(s, "A Self-Sustaining Fee Engine",
        Inches(0.55), Inches(0.55), Inches(12.0), Inches(0.7),
        fn=FH, fs=32, bold=True, color=C["white"])

    # Left — revenue table
    rect(s, Inches(0.4), Inches(1.4), Inches(6.2), Inches(5.6), fill=C["card"])
    section_label(s, "Fee Structure", Inches(0.65), Inches(1.6), C["amber"])

    rows = [
        ("Launch Fee",        "0.2 APT",   "Per token created",           C["amber"]),
        ("Pre-grad Trade",    "0.9% + 0.1%","Platform + creator",         C["blue"]),
        ("Post-grad Trade",   "0.05% + 0.2%+0.05%","Platform/Creator/LP",C["green"]),
        ("Graduation Fee",    "83 APT",    "60 platform + 23 creator",    C["purple"]),
    ]
    ry = Inches(2.0)
    for label, fee, note, color in rows:
        rect(s, Inches(0.6), ry, Inches(0.06), Inches(0.72), fill=color)
        txt(s, label, Inches(0.82), ry + Inches(0.05), Inches(2.2), Inches(0.38),
            fn=FB, fs=13, bold=True, color=C["light"])
        txt(s, fee,   Inches(3.2),  ry + Inches(0.05), Inches(1.5), Inches(0.38),
            fn=FH, fs=15, bold=True, color=color)
        txt(s, note,  Inches(0.82), ry + Inches(0.42), Inches(5.6), Inches(0.3),
            fn=FB, fs=10, color=C["muted"])
        ry += Inches(1.0)

    # Right — scenario modelling
    rect(s, Inches(6.9), Inches(1.4), Inches(6.0), Inches(5.6), fill=C["card2"])
    section_label(s, "Revenue Scenarios", Inches(7.1), Inches(1.6), C["green"])

    scenarios = [
        ("Conservative", "50 tokens/mo",   "~$2,800/mo",   "~$34K ARR",   C["muted"]),
        ("Base Case",    "300 tokens/mo",  "~$18,400/mo",  "~$220K ARR",  C["blue"]),
        ("Growth",       "1,000 tokens/mo","~$72,000/mo",  "~$864K ARR",  C["green"]),
        ("Viral (1%)",   "10K tokens/mo",  "~$680,000/mo", "~$8.2M ARR",  C["amber"]),
    ]
    ry = Inches(2.05)
    for name, vol, monthly, arr, color in scenarios:
        rect(s, Inches(7.0), ry, Inches(5.8), Inches(1.05), fill=C["card"])
        rect(s, Inches(7.0), ry, Inches(0.06), Inches(1.05), fill=color)
        txt(s, name,    Inches(7.2), ry + Inches(0.08), Inches(2.4), Inches(0.38),
            fn=FH, fs=14, bold=True, color=color)
        txt(s, vol,     Inches(7.2), ry + Inches(0.5),  Inches(2.4), Inches(0.38),
            fn=FB, fs=11, color=C["muted"])
        txt(s, monthly, Inches(9.5), ry + Inches(0.08), Inches(1.7), Inches(0.38),
            fn=FH, fs=15, bold=True, color=C["white"])
        txt(s, arr,     Inches(11.1),ry + Inches(0.08), Inches(1.6), Inches(0.38),
            fn=FB, fs=12, bold=True, color=color)
        ry += Inches(1.2)

    txt(s, "* Estimates based on 300 APT avg raise/token, APT at $8. Excludes post-graduation ongoing trade volume.",
        Inches(0.55), Inches(7.18), Inches(12.4), Inches(0.25),
        fn=FB, fs=9, color=C["muted"])
    return s


# ─── SLIDE 8 · SCALABILITY ARCHITECTURE ──────────────────────────────────────
def s08_scale(prs):
    s = blank(prs)
    bg(s, C["bg"])
    slide_chrome(s, "08", "Scalability")

    txt(s, "Built to Handle 100× Growth Without Rewrites",
        Inches(0.55), Inches(0.55), Inches(12.0), Inches(0.7),
        fn=FH, fs=32, bold=True, color=C["white"])
    txt(s, "3-phase architecture embedded in the codebase — not an afterthought.",
        Inches(0.55), Inches(1.25), Inches(10.0), Inches(0.38),
        fn=FB, fs=14, color=C["muted"])

    phases = [
        {
            "phase":  "PHASE 1",
            "status": "NOW",
            "title":  "Smart Polling",
            "color":  C["blue"],
            "points": [
                "Per-client HTTP polling to Geomi",
                "10-min TTL cache per user",
                "Request deduplication layer",
                "Visibility API stops polling",
                "when tab is hidden",
                "",
                "Capacity: ~500 concurrent users",
            ],
            "cost": "~$0.08 / 1K req",
        },
        {
            "phase":  "PHASE 2",
            "status": "NEXT",
            "title":  "WebSocket Broadcast",
            "color":  C["purple"],
            "points": [
                "1 server polls Geomi → pushes",
                "updates to all connected clients",
                "99.95% reduction in API calls",
                "Sub-second latency for all users",
                "Redis pub/sub for horizontal",
                "scaling across server instances",
                "",
                "Capacity: ~50,000 concurrent",
            ],
            "cost": "~$0.0004 / 1K req",
        },
        {
            "phase":  "PHASE 3",
            "status": "SCALE",
            "title":  "gRPC Streaming",
            "color":  C["green"],
            "points": [
                "Direct Aptos gRPC stream",
                "Zero polling — event-driven",
                "Real-time at blockchain speed",
                "Free data ingestion from node",
                "Unlimited concurrent users",
                "Horizontal microservice arch",
                "",
                "Capacity: Unlimited",
            ],
            "cost": "~$0 / event",
        },
    ]

    bx = Inches(0.4)
    bw = Inches(4.1)
    by = Inches(1.8)
    bh = Inches(5.3)

    for ph in phases:
        rect(s, bx, by, bw, bh, fill=C["card"])
        rect(s, bx, by, bw, Inches(0.07), fill=ph["color"])

        # Phase badge
        pill(s, ph["phase"], bx + Inches(0.22), by + Inches(0.2),
             fill=ph["color"], text_color=C["bg"])
        pill(s, ph["status"], bx + Inches(0.22) + Inches(1.1), by + Inches(0.2),
             fill=C["card2"], text_color=ph["color"])

        txt(s, ph["title"],
            bx + Inches(0.22), by + Inches(0.65), bw - Inches(0.35), Inches(0.55),
            fn=FH, fs=22, bold=True, color=C["white"])

        py2 = by + Inches(1.3)
        for pt in ph["points"]:
            if pt == "":
                py2 += Inches(0.15)
                continue
            col = C["light"] if not pt.startswith("  ") else C["muted"]
            txt(s, ("• " if not pt.startswith("  ") else "  ") + pt.strip(),
                bx + Inches(0.22), py2, bw - Inches(0.35), Inches(0.38),
                fn=FB, fs=12, color=col)
            py2 += Inches(0.41)

        # Cost badge bottom
        rect(s, bx + Inches(0.22), by + bh - Inches(0.55),
             bw - Inches(0.44), Inches(0.38), fill=ph["color"])
        txt(s, ph["cost"],
            bx + Inches(0.32), by + bh - Inches(0.54),
            bw - Inches(0.55), Inches(0.33),
            fn=FM, fs=11, bold=True, color=C["bg"])

        bx += Inches(4.45)

    # Arrow connectors
    for ax in [Inches(4.52), Inches(8.97)]:
        txt(s, "→", ax, Inches(4.2), Inches(0.4), Inches(0.5),
            fn=FH, fs=24, bold=True, color=C["blue"], align=PP_ALIGN.CENTER)

    return s


# ─── SLIDE 9 · COMPETITIVE MOAT ──────────────────────────────────────────────
def s09_moat(prs):
    s = blank(prs)
    bg(s, C["bg"])
    slide_chrome(s, "09", "Competitive Moat")

    txt(s, "Why MoveMint Wins — and Stays Won",
        Inches(0.55), Inches(0.55), Inches(12.0), Inches(0.7),
        fn=FH, fs=32, bold=True, color=C["white"])

    # Comparison table
    features = [
        "Aptos-native",
        "Bonding curve AMM",
        "Auto DEX graduation",
        "Locked liquidity (LP burn)",
        "Creator royalties on trades",
        "Real-time indexing",
        "Multi-wallet + Google OAuth",
        "Open source contract",
        "Scalable to gRPC",
    ]
    competitors = [
        ("MoveMint",    C["blue"],   [1,1,1,1,1,1,1,1,1]),
        ("Pump.fun",    C["muted"],  [0,1,1,1,0,1,0,0,0]),
        ("Cetus (Aptos)",C["muted"], [1,0,0,0,0,0,0,1,0]),
        ("DIY / Manual",C["muted"],  [1,0,0,0,0,0,0,0,0]),
    ]

    # Header row
    hx = Inches(5.2)
    cw = Inches(1.85)
    for comp, color, _ in competitors:
        rect(s, hx, Inches(1.4), cw, Inches(0.55), fill=C["card"])
        txt(s, comp, hx + Inches(0.12), Inches(1.45), cw - Inches(0.1), Inches(0.44),
            fn=FH, fs=13, bold=True, color=color, align=PP_ALIGN.CENTER)
        hx += cw + Inches(0.12)

    # Feature rows
    ry = Inches(2.05)
    for fi, feat in enumerate(features):
        row_fill = C["card"] if fi % 2 == 0 else C["card2"]
        rect(s, Inches(0.4), ry, Inches(4.6), Inches(0.56), fill=row_fill)
        txt(s, feat, Inches(0.55), ry + Inches(0.1), Inches(4.3), Inches(0.38),
            fn=FB, fs=13, color=C["light"])

        hx = Inches(5.2)
        for comp, color, checks in competitors:
            rect(s, hx, ry, cw, Inches(0.56), fill=row_fill)
            mark = "✓" if checks[fi] else "—"
            mc   = C["green"] if checks[fi] else C["divider"]
            txt(s, mark, hx, ry + Inches(0.08), cw, Inches(0.42),
                fn=FH, fs=18, bold=True, color=mc, align=PP_ALIGN.CENTER)
            hx += cw + Inches(0.12)

        ry += Inches(0.62)

    # Moat summary
    rect(s, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.38), fill=C["blue"])
    txt(s, "Network effects compound: every graduated token increases platform credibility and creator FOMO.",
        Inches(0.55), Inches(7.04), Inches(12.2), Inches(0.33),
        fn=FB, fs=12, bold=True, color=C["bg"])
    return s


# ─── SLIDE 10 · VISION ───────────────────────────────────────────────────────
def s10_vision(prs):
    s = blank(prs)
    bg(s, C["bg"])

    # Full-bleed gradient-style overlay (stacked rects, dark → purple tint)
    for i in range(30):
        alpha = i / 30
        r = int(0x0D + alpha * (0x1C - 0x0D))
        g = int(0x11 + alpha * (0x11 - 0x11))
        b = int(0x17 + alpha * (0x2E - 0x17))
        rect(s, 0, Inches(i * 0.25), W, Inches(0.26),
             fill=RGBColor(r, g, b))

    # Geometric accents
    rect(s, 0,            0,           Inches(0.06), H,   fill=C["purple"])
    rect(s, W-Inches(3),  0,           Inches(3),    Inches(0.07), fill=C["blue"])
    rect(s, Inches(0.55), Inches(2.85), Inches(6.0), Inches(0.05), fill=C["purple"])

    slide_chrome(s, "10", "Vision")

    txt(s, "The future of permissionless",
        Inches(0.55), Inches(1.0), Inches(12.0), Inches(1.0),
        fn=FH, fs=56, bold=True, color=C["white"])
    txt(s, "finance is being built now.",
        Inches(0.55), Inches(1.85), Inches(12.0), Inches(1.0),
        fn=FH, fs=56, bold=True, color=C["blue"])

    txt(s, "On Aptos. With Move. With MoveMint.",
        Inches(0.55), Inches(3.05), Inches(8.0), Inches(0.6),
        fn=FH, fs=26, bold=False, color=C["purple"])

    txt(s,
        "Our 18-month vision: become the canonical token launch infrastructure\n"
        "for Aptos — expanding to cross-chain graduation, creator DAOs,\n"
        "launchpad curation, and protocol-owned liquidity.",
        Inches(0.55), Inches(3.85), Inches(7.8), Inches(1.4),
        fn=FB, fs=16, color=C["muted"])

    # 3 vision KPI boxes
    vision_kpis = [
        (C["blue"],   "10,000+", "tokens launched\nby end of Year 1"),
        (C["purple"], "$50M+",   "total volume\nthrough bonding curves"),
        (C["green"],  "3",       "additional Aptos\nDEXs integrated"),
    ]
    vx = Inches(0.55)
    for color, big, label in vision_kpis:
        rect(s, vx, Inches(5.5), Inches(3.9), Inches(1.7), fill=C["card"])
        rect(s, vx, Inches(5.5), Inches(3.9), Inches(0.06), fill=color)
        txt(s, big,   vx + Inches(0.22), Inches(5.65),
            Inches(3.5), Inches(0.65),
            fn=FH, fs=36, bold=True, color=color)
        txt(s, label, vx + Inches(0.22), Inches(6.3),
            Inches(3.5), Inches(0.8),
            fn=FB, fs=13, color=C["muted"])
        vx += Inches(4.1)

    # CTA
    rect(s, Inches(0.55), Inches(7.1), Inches(8.5), Inches(0.32), fill=C["card2"])
    txt(s, "github.com/gking432/my-token-launcher-standard  ·  Ready for mainnet deployment",
        Inches(0.72), Inches(7.14), Inches(8.2), Inches(0.26),
        fn=FM, fs=10, color=C["muted"])

    txt(s, "MOVEMINT",
        Inches(9.5), Inches(6.55), Inches(3.5), Inches(0.75),
        fn=FH, fs=40, bold=True, color=C["purple"], align=PP_ALIGN.RIGHT)
    txt(s, "Aptos · Move · Open Source",
        Inches(9.5), Inches(7.15), Inches(3.5), Inches(0.3),
        fn=FB, fs=11, color=C["muted"], align=PP_ALIGN.RIGHT)
    return s


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN — BUILD & EXPORT
# ══════════════════════════════════════════════════════════════════════════════

def build():
    prs = prs_new()
    print("Building slides...")
    s01_hook(prs);    print("  ✓ Slide 01 — The Hook")
    s02_gap(prs);     print("  ✓ Slide 02 — The Gap")
    s03_thesis(prs);  print("  ✓ Slide 03 — Investment Thesis")
    s04_product(prs); print("  ✓ Slide 04 — The Innovation / Product")
    s05_tech(prs);    print("  ✓ Slide 05 — Technical Deep Dive")
    s06_journey(prs); print("  ✓ Slide 06 — User Journey")
    s07_revenue(prs); print("  ✓ Slide 07 — Market Strategy & Revenue")
    s08_scale(prs);   print("  ✓ Slide 08 — Scalability Architecture")
    s09_moat(prs);    print("  ✓ Slide 09 — Competitive Moat")
    s10_vision(prs);  print("  ✓ Slide 10 — Vision & CTA")

    out_dir  = os.path.dirname(os.path.abspath(__file__))
    pptx_path = os.path.join(out_dir, "MoveMint_Business_Case.pptx")
    pdf_path  = os.path.join(out_dir, "MoveMint_Business_Case.pdf")

    prs.save(pptx_path)
    print(f"\nSaved: {pptx_path}")

    # Convert to PDF via LibreOffice
    print("Converting to PDF via LibreOffice…")
    result = subprocess.run(
        ["libreoffice", "--headless", "--convert-to", "pdf",
         "--outdir", out_dir, pptx_path],
        capture_output=True, text=True
    )
    if result.returncode == 0:
        print(f"Saved: {pdf_path}")
    else:
        print("LibreOffice error:", result.stderr)

    print("\n✅ Done! Files created:")
    for f in [pptx_path, pdf_path]:
        if os.path.exists(f):
            size = os.path.getsize(f) / 1024
            print(f"   {f}  ({size:.1f} KB)")

    print("\n" + "═" * 60)
    print("  📋  5-MINUTE MANUAL CHECK CHECKLIST")
    print("═" * 60)
    print("  □  Slide 01 — Headline not clipped on right edge")
    print("  □  Slide 02 — Left/right panels don't overlap")
    print("  □  Slide 03 — All 3 pillars equal height")
    print("  □  Slide 04 — 6 cards in 2×3 grid, no overflow")
    print("  □  Slide 05 — Mono text readable inside code sections")
    print("  □  Slide 06 — 5 timeline cards evenly spaced")
    print("  □  Slide 07 — Revenue table columns aligned")
    print("  □  Slide 08 — Phase labels not cut off at bottom")
    print("  □  Slide 09 — Checkmark column alignment centered")
    print("  □  Slide 10 — Vision KPI boxes bottom-aligned")
    print("  □  ALL SLIDES — Left blue accent bar visible & consistent")
    print("  □  ALL SLIDES — Top-right purple accent bar visible")
    print("  □  ALL SLIDES — Slide numbers in top-left corner")
    print("═" * 60)

if __name__ == "__main__":
    build()
